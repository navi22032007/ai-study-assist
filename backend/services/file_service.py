import fitz  # PyMuPDF
import io
import base64
from typing import Tuple, List

def extract_text_from_pdf(file_data: bytes) -> Tuple[str, int]:
    """Extract text from PDF bytes. Returns (text, page_count)."""
    doc = fitz.open(stream=file_data, filetype="pdf")
    texts = []
    for page in doc:
        text = page.get_text("text")
        texts.append(text)
    doc.close()
    full_text = "\n\n".join(texts)
    return full_text, len(texts)

def extract_images_from_pdf(file_data: bytes) -> List[dict]:
    """
    Extract visual content from PDF pages.
    
    Strategy:
    1. First try extracting embedded raster images via get_images().
    2. If that finds nothing (common with vector diagrams, charts, academic figures),
       fall back to rendering each page as a high-res PNG screenshot.
       This captures EVERYTHING visible on the page regardless of how it's embedded.
    """
    doc = fitz.open(stream=file_data, filetype="pdf")
    images_result = []
    
    # ── Strategy 1: Extract embedded raster images ────────────────────────
    for page_num in range(len(doc)):
        page = doc[page_num]
        image_list = page.get_images(full=True)
        for img in image_list:
            xref = img[0]
            w = img[2]
            h = img[3]
            if w > 100 and h > 100:
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    ext = base_image["ext"].lower()
                    mime = "image/jpeg" if ext in ("jpg", "jpeg") else f"image/{ext}"
                    
                    if len(image_bytes) > 5 * 1024:  # > 5KB to skip icons
                        b64 = base64.b64encode(image_bytes).decode('utf-8')
                        images_result.append({
                            "data": b64,
                            "mime_type": mime,
                            "page": page_num + 1,
                            "source": "embedded"
                        })
                except Exception:
                    continue
            if len(images_result) >= 5:
                break
        if len(images_result) >= 5:
            break
    
    # ── Strategy 2: If no embedded images found, render pages as screenshots ─
    if len(images_result) == 0:
        # Render up to 5 pages at 2x resolution for clarity
        max_pages = min(len(doc), 5)
        for page_num in range(max_pages):
            page = doc[page_num]

            # Check if page has visual content beyond just text
            text_len = len(page.get_text("text").strip())
            has_drawings = len(page.get_drawings()) > 0
            has_images = len(page.get_images(full=True)) > 0

            # Skip pages that are primarily text with no visual elements
            if text_len > 200 and not has_drawings and not has_images:
                continue

            # Skip pages that are nearly empty
            if text_len < 10 and not has_drawings and not has_images:
                continue

            # Render page at 2x zoom for high quality
            mat = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("png")

            if len(img_bytes) > 10 * 1024:  # > 10KB (skip near-blank pages)
                b64 = base64.b64encode(img_bytes).decode('utf-8')
                images_result.append({
                    "data": b64,
                    "mime_type": "image/png",
                    "page": page_num + 1,
                    "source": "page_render"
                })
    
    doc.close()
    return images_result

def extract_text_from_docx(file_data: bytes) -> str:
    """Extract text from DOCX bytes."""
    try:
        import docx
        doc = docx.Document(io.BytesIO(file_data))
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        raise ValueError(f"Failed to extract text from DOCX: {e}")

def extract_text_from_pptx(file_data: bytes) -> str:
    """Extract text from PPTX bytes."""
    try:
        import pptx
        presentation = pptx.Presentation(io.BytesIO(file_data))
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        raise ValueError(f"Failed to extract text from PPTX: {e}")

def extract_text_from_txt(file_data: bytes) -> str:
    """Extract text from TXT bytes."""
    try:
        return file_data.decode("utf-8")
    except UnicodeDecodeError:
        return file_data.decode("latin-1", errors="replace")

def extract_text(file_data: bytes, file_type: str) -> str:
    """Extract text based on file type."""
    if file_type == "application/pdf":
        text, _ = extract_text_from_pdf(file_data)
        return text
    elif file_type in ["text/plain", "text/txt"]:
        return extract_text_from_txt(file_data)
    elif "wordprocessingml.document" in file_type or file_type == "application/msword":
        return extract_text_from_docx(file_data)
    elif "presentationml.presentation" in file_type or file_type == "application/vnd.ms-powerpoint":
        return extract_text_from_pptx(file_data)
    else:
        # Fallback to standard check or raise
        return extract_text_from_txt(file_data)

def validate_file(file_data: bytes, filename: str, content_type: str) -> Tuple[bool, str]:
    """Validate uploaded file. Returns (is_valid, error_message)."""
    MAX_SIZE = 50 * 1024 * 1024  # 50MB
    ALLOWED_EXTENSIONS = [".pdf", ".txt", ".docx", ".pptx", ".doc", ".ppt"]
    
    if len(file_data) > MAX_SIZE:
        return False, f"File size exceeds 50MB limit ({len(file_data) / 1024 / 1024:.1f}MB)"
    
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in ALLOWED_EXTENSIONS:
        return False, f"Unsupported file type. Allowed: PDF, TXT, DOCX, PPTX."
    
    return True, ""

